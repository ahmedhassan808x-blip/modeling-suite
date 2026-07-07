"""
Deck builder — navy & steel, built from blank layouts (no default-template
look). 16:9, title slide with a navy band, content slides with a title rule,
footer with source line and page numbers.

Model report modules compose slides from these primitives; the numbers on
every slide come from the recalculated workbook, so deck and model can never
disagree.
"""

from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from shared.theme import FONT, GRAY, INK, MIST, NAVY, SLATE, STEEL

C_NAVY, C_STEEL, C_SLATE = (RGBColor.from_string(c) for c in (NAVY, STEEL, SLATE))
C_INK, C_GRAY, C_MIST = (RGBColor.from_string(c) for c in (INK, GRAY, MIST))
C_WHITE = RGBColor.from_string("FFFFFF")

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def _text(frame, runs, size=12, color=C_INK, bold=False, align=PP_ALIGN.LEFT,
          space_after=6):
    """runs: str or list of (text, overrides-dict) paragraphs."""
    if isinstance(runs, str):
        runs = [(runs, {})]
    frame.word_wrap = True
    for i, (txt, ov) in enumerate(runs):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = ov.get("align", align)
        p.space_after = Pt(ov.get("space_after", space_after))
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = FONT
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)


class Deck:
    def __init__(self, footer: str):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = SLIDE_W, SLIDE_H
        self.blank = self.prs.slide_layouts[6]
        self.footer = footer
        self._n = 0

    # ---- primitives -------------------------------------------------------

    def _box(self, slide, x, y, w, h, fill=None, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sh.fill.solid() if fill else sh.fill.background()
        if fill:
            sh.fill.fore_color.rgb = fill
        sh.line.fill.background()
        if line:
            sh.line.fill.solid()
            sh.line.fill.fore_color.rgb = line
        sh.shadow.inherit = False
        # Strip the theme style reference — it carries an effectRef whose
        # drop shadow renderers apply despite shadow.inherit = False.
        style = sh._element.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/"
            "presentationml}style")
        if style is None:
            style = next((el for el in sh._element
                          if el.tag.endswith("}style")), None)
        if style is not None:
            sh._element.remove(style)
        return sh

    def _footer(self, slide):
        self._n += 1
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.08),
                                      Inches(11.0), Inches(0.35))
        _text(tb.text_frame, self.footer, size=8, color=C_GRAY)
        pn = slide.shapes.add_textbox(Inches(12.4), Inches(7.08),
                                      Inches(0.6), Inches(0.35))
        _text(pn.text_frame, str(self._n), size=8, color=C_GRAY,
              align=PP_ALIGN.RIGHT)

    # ---- slides -----------------------------------------------------------

    def title_slide(self, title, subtitle, tagline=""):
        s = self.prs.slides.add_slide(self.blank)
        self._box(s, 0, 0, SLIDE_W, Inches(0.18), fill=C_STEEL)
        self._box(s, 0, Inches(2.35), SLIDE_W, Inches(2.1), fill=C_NAVY)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(2.62), Inches(11.5),
                                  Inches(1.6))
        _text(tb.text_frame,
              [(title, dict(size=34, bold=True, color=C_WHITE)),
               (subtitle, dict(size=16, color=RGBColor.from_string(STEEL)))])
        tb2 = s.shapes.add_textbox(Inches(0.9), Inches(5.1), Inches(11.5),
                                   Inches(1.0))
        _text(tb2.text_frame,
              [(date.today().strftime("%B %d, %Y"), dict(size=12,
                                                         color=C_SLATE)),
               (tagline, dict(size=10, color=C_GRAY))])
        return s

    def content_slide(self, title):
        s = self.prs.slides.add_slide(self.blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3),
                                  Inches(0.65))
        _text(tb.text_frame, title, size=22, bold=True, color=C_NAVY)
        self._box(s, Inches(0.5), Inches(1.0), Inches(12.33), Emu(28575),
                  fill=C_STEEL)  # 0.03" rule under the title
        self._footer(s)
        return s

    def add_bullets(self, slide, items, x, y, w, h, size=13):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w),
                                      Inches(h))
        runs = []
        for it in items:
            if isinstance(it, tuple):  # (text, indent-level-ignored/emphasis)
                runs.append((f"— {it[0]}", dict(size=size - 2, color=C_SLATE)))
            else:
                runs.append((f"•  {it}", dict(size=size, space_after=10)))
        _text(tb.text_frame, runs)
        return tb

    def add_image(self, slide, png, x, y, w):
        return slide.shapes.add_picture(str(png), Inches(x), Inches(y),
                                        width=Inches(w))

    def add_table(self, slide, rows, x, y, w, size=10, first_col_w=None):
        """rows[0] is the header. Navy header, banded body, right-aligned data."""
        n_r, n_c = len(rows), len(rows[0])
        shape = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                                       Inches(w), Inches(0.3 * n_r))
        tbl = shape.table
        tbl.first_row = False
        tbl.horz_banding = False
        if first_col_w:
            tbl.columns[0].width = Inches(first_col_w)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    C_NAVY if i == 0 else C_MIST if i % 2 == 0 else C_WHITE)
                _text(cell.text_frame, str(val), size=size,
                      bold=i == 0 or j == 0,
                      color=C_WHITE if i == 0 else C_INK,
                      align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT,
                      space_after=0)
                cell.margin_top = cell.margin_bottom = Pt(3)
        return tbl

    def save(self, path):
        self.prs.save(str(path))
        return path
