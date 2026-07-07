"""
Presentation layer: all four export formats produced per model, decks
structurally sound (slide count, title, tables), numbers sourced from the
recalculated workbook, and the pipeline refuses workbooks with formula errors.
"""

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from dcf.model_builder import build_dcf_model
from dcf.report import build_reports as dcf_reports, extract_dcf
from three_statement.assumptions import Assumptions
from three_statement.model_builder import build_model
from three_statement.report import build_reports as ts_reports, extract
from tests.test_dcf import FAKE_PEERS


@pytest.fixture(scope="module")
def ts_outputs(simpleco_module, tmp_path_factory):
    tmp = tmp_path_factory.mktemp("ts_report")
    xlsx = tmp / "SMPL_3stmt.xlsx"
    build_model(simpleco_module, Assumptions.derive(simpleco_module), xlsx)
    return xlsx, ts_reports(xlsx)


@pytest.fixture(scope="module")
def dcf_outputs(simpleco_module, tmp_path_factory):
    tmp = tmp_path_factory.mktemp("dcf_report")
    xlsx = tmp / "SMPL_dcf.xlsx"
    build_dcf_model(simpleco_module, Assumptions.derive(simpleco_module),
                    xlsx, peers=FAKE_PEERS)
    return xlsx, dcf_reports(xlsx)


def _check_outputs(outputs, min_slides):
    for kind in ("pptx", "docx", "pdf"):
        assert outputs[kind].exists(), f"{kind} missing"
        assert outputs[kind].stat().st_size > 10_000, f"{kind} suspiciously small"
    assert outputs["pdf"].read_bytes()[:5] == b"%PDF-"
    prs = Presentation(outputs["pptx"])
    assert len(prs.slides) >= min_slides
    return prs


def test_three_statement_reports(ts_outputs):
    xlsx, outputs = ts_outputs
    prs = _check_outputs(outputs, min_slides=6)
    title_texts = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes
                           if sh.has_text_frame)
    assert "SimpleCo" in title_texts
    doc = Document(outputs["docx"])
    body = "\n".join(p.text for p in doc.paragraphs)
    assert "Executive summary" in body and "Limitations" in body


def test_dcf_reports(dcf_outputs):
    xlsx, outputs = dcf_outputs
    prs = _check_outputs(outputs, min_slides=6)
    all_text = " ".join(sh.text_frame.text for s in prs.slides
                        for sh in s.shapes if sh.has_text_frame)
    assert "DCF" in all_text and "WACC" in all_text
    # Deck numbers must equal the recalculated workbook's numbers
    d = extract_dcf(xlsx)
    assert f"${d['dcf']['ps']:,.2f}" in all_text


def test_extract_matches_workbook(ts_outputs, simpleco_module):
    xlsx, _ = ts_outputs
    d = extract(xlsx)
    assert d["ticker"] == "SMPL"
    assert d["checks"] == "PASS"
    assert d["revenue"][:3] == pytest.approx([1000, 1100, 1210])
    assert len(d["revenue"]) == 8


def test_dcf_extract_value_composition(dcf_outputs):
    xlsx, _ = dcf_outputs
    d = extract_dcf(xlsx)
    share = d["dcf"]["pv_tv"] / (d["dcf"]["sum_pv"] + d["dcf"]["pv_tv"])
    assert 0.3 < share < 0.95          # TV dominates but isn't everything
    assert len(d["field"]) >= 3        # football field has real rows
    assert len(d["ufcf"]) == 5


def test_report_refuses_broken_workbook(simpleco_module, tmp_path):
    xlsx = tmp_path / "broken.xlsx"
    build_model(simpleco_module, Assumptions.derive(simpleco_module), xlsx)
    wb = load_workbook(xlsx)
    wb["IS"]["Z99"] = "=1/0"
    wb.save(xlsx)
    with pytest.raises(RuntimeError, match="formula errors"):
        extract(xlsx)
